# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
PRIMARY_BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK_BLUE = RGBColor(0x0D, 0x33, 0x8A)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xFF)
ACCENT_GOLD = RGBColor(0xFF, 0xB3, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)

W = Cm(33.867)
H = Cm(19.05)
TOTAL_SLIDES = 14

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def add_bg_rect(slide, color, left=0, top=0, width=None, height=None):
    w = width or W
    h = height or H
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei', anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(tf, text, font_size=16, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei', space_before=Pt(4), space_after=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_bullet(tf, text, font_size=15, color=BLACK, font_name='Microsoft YaHei', level=0, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.level = level
    p.font.bold = bold
    p.space_before = Pt(3)
    p.space_after = Pt(3)
    return p

def add_title_bar(slide, text):
    add_bg_rect(slide, PRIMARY_BLUE, 0, 0, W, Cm(2.2))
    add_textbox(slide, Cm(1.5), Cm(0.3), Cm(30), Cm(1.6), text, font_size=26, bold=True, color=WHITE)
    add_bg_rect(slide, ACCENT_GOLD, 0, Cm(2.2), W, Cm(0.08))

def add_decoration_line(slide, left, top, width=Cm(3), height=Cm(0.06)):
    add_bg_rect(slide, ACCENT_GOLD, left, top, width, height)

def add_page_number(slide, num):
    add_textbox(slide, W - Cm(3), H - Cm(0.8), Cm(2.5), Cm(0.7), f'{num}/{TOTAL_SLIDES}', font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

# ===== Slide 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, DARK_BLUE)
add_bg_rect(slide, PRIMARY_BLUE, 0, 0, W, Cm(0.3))
add_textbox(slide, Cm(2), Cm(1.5), Cm(30), Cm(1.5), '浙江水利水电学院', font_size=20, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Cm(2), Cm(2.8), Cm(30), Cm(1.2), '本科毕业设计答辩', font_size=18, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Cm(2), Cm(5), Cm(30), Cm(3), '保温时间对快速加热IF复合板\n组织和性能的研究', font_size=38, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_decoration_line(slide, Cm(12), Cm(8.2), Cm(10))

info_items = [
    ('学    院：', '机械工程学院'),
    ('专    业：', '材料成型及控制工程'),
    ('学    号：', '2348003113'),
    ('学生姓名：', '张新'),
    ('指导教师：', '唐恩'),
]
for i, (label, value) in enumerate(info_items):
    y = Cm(9.5 + i * 1.1)
    add_textbox(slide, Cm(10), y, Cm(6), Cm(1), label, font_size=16, color=LIGHT_BLUE, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, Cm(16), y, Cm(8), Cm(1), value, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

add_textbox(slide, Cm(2), H - Cm(1.5), Cm(30), Cm(1), '2025届  ·  答辩时间：2025年5月', font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 1)

# ===== Slide 2: TOC =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '目  录')
toc_items = [
    ('01', '研究背景与意义'),
    ('02', '国内外研究现状'),
    ('03', '研究内容与创新点'),
    ('04', '实验材料与方法'),
    ('05', '结果与讨论'),
    ('06', '结论与展望'),
]
for i, (num, title) in enumerate(toc_items):
    col = i % 2
    row = i // 2
    x = Cm(2 + col * 16)
    y = Cm(4 + row * 4)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Cm(1.8), Cm(1.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    add_textbox(slide, x + Cm(2.5), y + Cm(0.2), Cm(12), Cm(1.4), title, font_size=20, bold=True, color=DARK_BLUE)
add_page_number(slide, 2)

# ===== Slide 3: Background =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '一、研究背景与意义')
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(31), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

add_para(tf, 'IF钢与铜复合板', font_size=20, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, 'IF钢（无间隙原子钢）：添加Ti、Nb固定C、N原子，具有优异的深冲性能和非时效性')
add_bullet(tf, '铜：出色的导电性、导热性及良好的耐腐蚀性')
add_bullet(tf, '复合后兼具高强度、易加工性和优良导电导热特性')
add_bullet(tf, '应用领域：电子设备散热部件、汽车发动机关键组件等')

add_para(tf, '', font_size=6)
add_para(tf, '研究必要性', font_size=20, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '快速加热工艺可显著影响复合板组织和性能，但保温时间优化研究不足')
add_bullet(tf, '缺乏多温度梯度下保温时间对IF复合板性能的系统研究')
add_bullet(tf, '保温时间与界面结合、晶粒细化、力学性能的定量关系尚不明确')
add_bullet(tf, '本研究填补该领域空白，为生产工艺优化提供理论依据')
add_page_number(slide, 3)

# ===== Slide 4: Research Status =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '二、国内外研究现状')
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(31), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

add_para(tf, '异种金属复合界面扩散行为', font_size=18, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '钛/钢复合板：750-800℃扩散退火时界面剪切强度最高可达210MPa')
add_bullet(tf, '过高或过低温度均会导致界面剪切强度下降')

add_para(tf, '', font_size=6)
add_para(tf, '热处理工艺对IF钢晶粒度的影响', font_size=18, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '加热速度、温度、保温时间和冷却速度均对晶粒度产生显著影响')
add_bullet(tf, '较高加热温度和较长保温时间通常导致晶粒长大')
add_bullet(tf, 'CSP线数据：碳含量50ppm以内时性能最优（r>2.0, n>0.22）')

add_para(tf, '', font_size=6)
add_para(tf, '保温时间与界面剪切强度的相关性', font_size=18, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '适当延长保温时间促进元素扩散，增强界面结合')
add_bullet(tf, '过长保温时间可能生成脆性相，导致强度下降')
add_bullet(tf, '存在最佳保温时间窗口使界面性能最优')
add_page_number(slide, 4)

# ===== Slide 5: Content & Innovation =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '三、研究内容与创新点')
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(31), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

add_para(tf, '研究内容', font_size=20, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '系统探究多温度梯度（830℃、850℃、870℃）下保温时间对IF复合板的影响')
add_bullet(tf, '保温时间设置：0s、4s、8s、12s，加热后统一水冷处理')
add_bullet(tf, '金相分析 + 晶粒度测定 + 硬度测试 → 确定最优参数')
add_bullet(tf, '最优参数下进行拉伸试验 + 界面SEM检测')

add_para(tf, '', font_size=6)
add_para(tf, '技术路线', font_size=20, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, 'IF复合板 → 到温入炉加热(830/850/870℃) → 保温(0/4/8/12s) → 水冷淬火')
add_bullet(tf, '→ 金相分析 & 晶粒度测定 & 硬度测试 → 确定最优参数')
add_bullet(tf, '→ 拉伸试验 + 界面SEM断口分析')

add_para(tf, '', font_size=6)
add_para(tf, '创新点', font_size=20, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '多温度梯度（830/850/870℃）系统对比保温时间效应')
add_bullet(tf, '快速加热工艺与常规工艺对比，揭示快速加热对界面结合和晶粒细化的影响机制')
add_bullet(tf, '微观（SEM断口分析）与宏观（拉伸、硬度）相结合的综合研究方法')
add_page_number(slide, 5)

# ===== Slide 6: Experimental =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '四、实验材料与方法')

txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(15), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True
add_para(tf, '实验材料', font_size=20, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, 'IF钢：0.002%C, Ti 0.04%, Nb 0.03%')
add_bullet(tf, '纯铜：纯度99.9%')
add_bullet(tf, '复合板厚度：1.5mm')

add_para(tf, '', font_size=6)
add_para(tf, '加热设备', font_size=20, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '常规电阻炉（硅碳棒加热）')
add_bullet(tf, '温度范围：室温~1200℃')
add_bullet(tf, '控温精度：±1℃，升温速率10℃/min')

txBox2 = slide.shapes.add_textbox(Cm(17), Cm(3), Cm(16), Cm(14))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_para(tf2, '实验方案', font_size=20, bold=True, color=PRIMARY_BLUE)
add_bullet(tf2, '加热温度：830℃ / 850℃ / 870℃')
add_bullet(tf2, '保温时间：0s / 4s / 8s / 12s')
add_bullet(tf2, '冷却方式：水冷淬火（20℃去离子水）')
add_bullet(tf2, '到温入炉，电子计时器精确计时')

add_para(tf2, '', font_size=6)
add_para(tf2, '性能测试', font_size=20, bold=True, color=PRIMARY_BLUE)
add_bullet(tf2, '金相观察（蔡司Axioscope5, 500×）')
add_bullet(tf2, '晶粒度评级（ASTM E112, 截点法）')
add_bullet(tf2, '显微硬度（HV-1000, 25g/5s）')
add_bullet(tf2, '拉伸试验（CMT5105, 1mm/min）')
add_bullet(tf2, 'SEM断口分析（喷金处理）')
add_page_number(slide, 6)

# ===== Slide 7: Results - Metallography & Grain Size =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '五、结果与讨论——金相组织与晶粒度')
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(31), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

add_para(tf, '金相组织观察', font_size=20, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '830℃：保温4s组织细小均匀（19.66μm）；8s略细（18.43μm）；12s略粗（21.78μm）')
add_bullet(tf, '850℃：随保温时间延长，晶粒持续长大（22.29→24.81→26.85μm）')
add_bullet(tf, '870℃：晶粒长大更显著（25.82→26.83→27.25μm）')

add_para(tf, '', font_size=6)
add_para(tf, '晶粒度数据（ASTM E112, 截点法）', font_size=18, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '原件：8.6级 | 平均21.22μm')
add_bullet(tf, '830℃-4s：8.7级 | 19.66μm  ★ 晶粒细化')
add_bullet(tf, '830℃-8s：8.8级 | 18.43μm  ★ 最细')
add_bullet(tf, '850℃-8s：8.4级 | 24.81μm  → 综合优化')
add_bullet(tf, '870℃-4s：8.3级 | 25.82μm  → 高温优化')
add_bullet(tf, '870℃-12s：8.2级 | 27.25μm  ↑ 晶粒粗化')
add_page_number(slide, 7)

# ===== Slide 8: Hardness =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '五、结果与讨论——硬度分析')
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(31), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

add_para(tf, '不同温度下硬度随保温时间的变化趋势', font_size=20, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '总体趋势：随保温时间延长，硬度总体呈下降趋势', bold=True)
add_para(tf, '', font_size=4)
add_bullet(tf, '830℃：4s最高139.62HV → 8s 134.08HV → 12s 132.75HV')
add_bullet(tf, '850℃：4s 137.17HV → 8s 138.69HV（略升）→ 12s 131.33HV（显著降）')
add_bullet(tf, '870℃：4s最高138.45HV → 8s 131.89HV → 12s 130.80HV')

add_para(tf, '', font_size=6)
add_para(tf, '晶粒度与硬度的响应规律（Hall-Petch关系）', font_size=18, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '晶粒细化强化：热处理后硬度普遍高于原始（130.52HV），再结晶形成细小等轴晶')
add_bullet(tf, '晶粒长大软化：同一温度下随保温时间延长，晶粒增大→硬度下降')
add_bullet(tf, '830℃晶粒细化显著→硬度升高；870℃晶粒长大加剧→硬度降低')
add_bullet(tf, '位错密度、第二相析出、界面扩散等也与晶粒尺寸产生交互作用')
add_page_number(slide, 8)

# ===== Slide 9: Tensile =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '五、结果与讨论——拉伸性能')
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(31), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

add_para(tf, '最优参数下的拉伸试验结果', font_size=20, bold=True, color=PRIMARY_BLUE)
add_para(tf, '', font_size=4)
add_para(tf, '830℃/4s：抗拉300~310MPa  伸长率47~49.5%  弹性模量73004~78179MPa', font_size=14, color=DARK_GRAY)
add_para(tf, '850℃/8s：抗拉305MPa      伸长率60.5~64%  弹性模量79596~80177MPa', font_size=14, color=DARK_GRAY)
add_para(tf, '870℃/4s：抗拉305MPa      伸长率62.5~65.5%  弹性模量73980~76933MPa', font_size=14, color=DARK_GRAY)

add_para(tf, '', font_size=6)
add_para(tf, '拉伸性能分析', font_size=18, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '830℃保温4s：强度较高（300~310MPa），但伸长率较低（~48%），塑性相对不足')
add_bullet(tf, '850℃保温8s：强度适中（305MPa），伸长率显著提高（~62%），综合性能优异')
add_bullet(tf, '870℃保温4s：强度与850℃相近，伸长率最高（~64%），塑性最优')
add_para(tf, '', font_size=4)
add_para(tf, '结论：850℃/8s和870℃/4s为较优拉伸工艺参数，兼顾强度和塑性', font_size=16, bold=True, color=PRIMARY_BLUE)
add_page_number(slide, 9)

# ===== Slide 10: SEM =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '五、结果与讨论——界面断口SEM分析')
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(31), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

add_para(tf, '830℃保温4s断口特征', font_size=18, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '可见大量韧窝（韧性断裂特征），存在撕裂棱（含脆性成分）')
add_bullet(tf, '存在明显界面分离现象 → IF钢与铜结合强度较弱')
add_bullet(tf, '断裂机制：微孔洞形核→长大聚合→韧性断裂+界面脱粘')

add_para(tf, '', font_size=4)
add_para(tf, '850℃保温8s断口特征（最佳）', font_size=18, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '韧窝尺寸更大、形状更规则、分布更均匀 → 塑性更好')
add_bullet(tf, '撕裂棱明显减少，界面分离现象显著减少')
add_bullet(tf, '断裂机制以韧性断裂为主 → 界面结合强度提高')

add_para(tf, '', font_size=4)
add_para(tf, '870℃保温4s断口特征', font_size=18, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '韧窝尺寸较大但分布不均，存在较大孔洞')
add_bullet(tf, '界面分离现象略有增加，脆性成分有所增加')
add_bullet(tf, '高温导致界面扩散加剧，但结合强度有所下降')
add_page_number(slide, 10)

# ===== Slide 11: Optimization =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '五、结果与讨论——最优工艺参数综合对比')
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(31), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

add_para(tf, '三组最优参数综合对比', font_size=20, bold=True, color=PRIMARY_BLUE)
add_para(tf, '', font_size=4)
add_para(tf, '830℃/4s    8.7级    139.62HV    300~310MPa    48.25%    硬度高，塑性一般', font_size=14, color=DARK_GRAY)
add_para(tf, '850℃/8s    8.4级    138.69HV    305MPa        62.25%    综合性能优异', font_size=14, color=PRIMARY_BLUE, bold=True)
add_para(tf, '870℃/4s    8.3级    138.45HV    305MPa        64.00%    塑性最优', font_size=14, color=DARK_GRAY)

add_para(tf, '', font_size=6)
add_para(tf, '与常规退火工艺对比', font_size=18, bold=True, color=PRIMARY_BLUE)
add_bullet(tf, '常规退火工艺：加热速度慢，保温时间长（数分钟到数十分钟）')
add_bullet(tf, '常规工艺下：抗拉强度约400MPa，延伸率约25%，平均晶粒尺寸约15um')
add_bullet(tf, '本工艺（快速加热+优化保温）：强度约305MPa，延伸率高达64%，界面结合良好')
add_para(tf, '', font_size=4)
add_bullet(tf, '结论：快速加热工艺在优化保温时间后，能够显著改善IF复合板的综合力学性能', bold=True)
add_page_number(slide, 11)

# ===== Slide 12: Conclusion =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '六、研究结论')
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(31), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

add_para(tf, '主要研究结论', font_size=22, bold=True, color=PRIMARY_BLUE)
add_para(tf, '', font_size=6)

add_para(tf, '1. 成功建立了保温时间-温度-组织性能响应模型', font_size=16, bold=True, color=DARK_BLUE)
add_bullet(tf, '830℃：晶粒细化显著，4s时平均晶粒尺寸19.66um，硬度最高139.62HV')
add_bullet(tf, '850℃：晶粒长大加速，8s时综合性能最优（伸长率62.25%，硬度138.69HV）')
add_bullet(tf, '870℃：晶粒粗化明显，4s时塑性最优（伸长率64%）')

add_para(tf, '', font_size=4)
add_para(tf, '2. 确定了三组最优工艺参数', font_size=16, bold=True, color=DARK_BLUE)
add_bullet(tf, '830℃保温4s / 850℃保温8s / 870℃保温4s')
add_bullet(tf, '三者均可实现良好的界面结合和力学性能')

add_para(tf, '', font_size=4)
add_para(tf, '3. 界面分析验证了良好的结合质量', font_size=16, bold=True, color=DARK_BLUE)
add_bullet(tf, '拉伸后界面处未出现脱粘或分层现象')
add_bullet(tf, 'SEM显示以韧性断裂为主，证实界面结合质量良好')
add_page_number(slide, 12)

# ===== Slide 13: Outlook =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, '七、研究展望')
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(31), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

add_para(tf, '研究展望', font_size=22, bold=True, color=PRIMARY_BLUE)
add_para(tf, '', font_size=6)

add_para(tf, '1. 长期服役性能评估', font_size=18, bold=True, color=DARK_BLUE)
add_bullet(tf, '开展疲劳寿命试验，研究循环载荷下的裂纹萌生和扩展机制')
add_bullet(tf, '进行耐腐蚀性能研究，评估不同腐蚀介质中的服役行为')

add_para(tf, '', font_size=4)
add_para(tf, '2. 工业化应用推广', font_size=18, bold=True, color=DARK_BLUE)
add_bullet(tf, '开发适合工业化生产的快速加热设备')
add_bullet(tf, '优化生产工艺，提高生产效率和产品质量稳定性')
add_bullet(tf, '开展中试试验，验证工艺在工业化生产中的可行性')

add_para(tf, '', font_size=4)
add_para(tf, '3. 跨学科交叉融合', font_size=18, bold=True, color=DARK_BLUE)
add_bullet(tf, '材料科学 + 机械工程 + 自动化控制协同创新')
add_bullet(tf, '推动快速加热工艺的规模化应用')
add_page_number(slide, 13)

# ===== Slide 14: Thanks =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, DARK_BLUE)
add_bg_rect(slide, PRIMARY_BLUE, 0, 0, W, Cm(0.3))
add_textbox(slide, Cm(2), Cm(3), Cm(30), Cm(2), '致  谢', font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_decoration_line(slide, Cm(13), Cm(5.5), Cm(8))
add_textbox(slide, Cm(2), Cm(7), Cm(30), Cm(3), '感谢指导教师唐恩老师的悉心指导\n感谢实验室师兄师姐们的帮助\n感谢浙江水利水电学院机械工程学院的培养', font_size=18, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Cm(2), Cm(12), Cm(30), Cm(2), '请各位评委老师批评指正！', font_size=24, bold=True, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Cm(2), H - Cm(1.5), Cm(30), Cm(1), '张新  ·  机械工程学院  ·  材料成型及控制工程', font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 14)

# Save
output_path = r'D:\桌面\abc\张新_毕业设计答辩PPT.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
