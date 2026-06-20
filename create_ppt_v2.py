# -*- coding: utf-8 -*-
"""毕业设计答辩PPT v2 - 专业设计版"""
 
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')
 
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
 
# === Color Palette ===
CP = RGBColor(0x1B, 0x3A, 0x6B)      # 主色深蓝
CS = RGBColor(0x2D, 0x6D, 0xCC)      # 辅色中蓝
CA = RGBColor(0x00, 0xA8, 0xE8)      # 强调亮蓝
CG = RGBColor(0xF5, 0xA6, 0x23)      # 金色
CW = RGBColor(0xFF, 0xFF, 0xFF)      # 白色
CB = RGBColor(0x2C, 0x3E, 0x50)      # 深灰黑
CD = RGBColor(0x4A, 0x55, 0x66)      # 深灰
CH = RGBColor(0xE8, 0xF0, 0xFE)      # 浅蓝高亮
CTH = RGBColor(0x1B, 0x3A, 0x6B)     # 表头
CTA = RGBColor(0xEB, 0xF1, 0xFB)     # 表格交替色
CMG = RGBColor(0x7F, 0x8C, 0x9A)     # 中灰
 
W = Cm(33.867)
H = Cm(19.05)
TOTAL = 14
 
prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]
 
def rect(sl, c, l, t, w, h):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh
 
def rrect(sl, c, l, t, w, h):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh
 
def tb(sl, l, t, w, h, txt, sz=18, b=False, c=CB, a=PP_ALIGN.LEFT, anc=MSO_ANCHOR.TOP):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True; tf.auto_size = None
    try: tf.vertical_anchor = anc
    except: pass
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz)
    p.font.bold = b; p.font.color.rgb = c; p.font.name = 'Microsoft YaHei'; p.alignment = a
    return tf
 
def para(tf, txt, sz=16, b=False, c=CB, a=PP_ALIGN.LEFT):
    p = tf.add_paragraph(); p.text = txt; p.font.size = Pt(sz)
    p.font.bold = b; p.font.color.rgb = c; p.font.name = 'Microsoft YaHei'; p.alignment = a
    p.space_before = Pt(3); p.space_after = Pt(3); return p
 
def blt(tf, txt, sz=15, c=CB, b=False):
    p = tf.add_paragraph(); p.text = txt; p.font.size = Pt(sz)
    p.font.color.rgb = c; p.font.name = 'Microsoft YaHei'; p.font.bold = b
    p.space_before = Pt(3); p.space_after = Pt(3); return p
 
def tbar(sl, txt):
    rect(sl, CP, 0, 0, W, Cm(2.5))
    rect(sl, CA, 0, Cm(2.5), W, Cm(0.12))
    tb(sl, Cm(2), Cm(0.3), Cm(30), Cm(1.2), txt, sz=28, b=True, c=CW)
 
def pgn(sl, n):
    tb(sl, W-Cm(3), H-Cm(0.8), Cm(2.5), Cm(0.6), f'{n}/{TOTAL}', sz=9, c=CMG, a=PP_ALIGN.RIGHT)
 
def sdeco(sl):
    rect(sl, CS, Cm(0.6), Cm(3), Cm(0.08), Cm(14))
 
def mtbl(sl, data, l, t, w=None, cw=None):
    """Create a table from data list. First row is header."""
    rows, cols = len(data), len(data[0])
    sh = sl.shapes.add_table(rows, cols, l, t, w or Cm(31), None)
    tbl = sh.table
    # set column widths
    if cw:
        for i, c in enumerate(cw):
            if i < cols: tbl.columns[i].width = c
    # fill cells
    for ri, rd in enumerate(data):
        h = Pt(18) if ri == 0 else Pt(36)
        for ci, ct in enumerate(rd):
            cell = tbl.cell(ri, ci)
            cell.text = str(ct)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12) if ri > 0 else Pt(12)
                p.font.name = 'Microsoft YaHei'
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = (ri == 0)
                p.font.color.rgb = CW if ri == 0 else CB
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = CTH
            elif ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = CTA
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tbl
 
def oval(sl, l, t, s, c, txt, sz=14):
    sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, s, s)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    tf = sh.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz)
    p.font.bold = True; p.font.color.rgb = CW; p.alignment = PP_ALIGN.CENTER
 
# ========== Slide 1: Title ==========
sl = prs.slides.add_slide(blank)
rect(sl, CP, 0, 0, W, H)
rect(sl, CA, 0, 0, W, Cm(0.3))
tb(sl, Cm(2), Cm(2), Cm(30), Cm(1.5), '浙江水利水电学院', sz=22, c=CH, a=PP_ALIGN.CENTER)
tb(sl, Cm(2), Cm(3.3), Cm(30), Cm(1), '本科毕业设计答辩', sz=17, c=CG, a=PP_ALIGN.CENTER)
rect(sl, CG, Cm(10), Cm(5), Cm(14), Cm(0.06))
tb(sl, Cm(3.5), Cm(6), Cm(27), Cm(4), '保温时间对快速加热IF复合板\n组织和性能的研究',
   sz=40, b=True, c=CW, a=PP_ALIGN.CENTER)
rect(sl, CG, Cm(12), Cm(9.5), Cm(10), Cm(0.04))
info = [('学    院：', '机械工程学院'), ('专    业：', '材料成型及控制工程'),
        ('学    号：', '2348003113'), ('学生姓名：', '张新'), ('指导教师：', '唐恩')]
for i, (l, v) in enumerate(info):
    y = Cm(10.8 + i*1.05)
    tb(sl, Cm(9), y, Cm(6), Cm(0.8), l, sz=15, c=CH, a=PP_ALIGN.RIGHT)
    tb(sl, Cm(15), y, Cm(8), Cm(0.8), v, sz=15, b=True, c=CW)
tb(sl, Cm(2), H-Cm(1.2), Cm(30), Cm(0.8), '2025届  |  答辩时间：2025年5月', sz=11, c=CMG, a=PP_ALIGN.CENTER)
pgn(sl, 1)
 
# ========== Slide 2: TOC ==========
sl = prs.slides.add_slide(blank)
tbar(sl, '目  录')
toc = [('01', '研究背景与意义'), ('02', '国内外研究现状'), ('03', '研究内容与创新点'),
       ('04', '实验材料与方法'), ('05', '结果与讨论'), ('06', '结论与展望')]
for i, (n, t) in enumerate(toc):
    col, row = i%2, i//2
    x, y = Cm(2+col*16.5), Cm(4+row*4.2)
    oval(sl, x, y, Cm(1.6), CP, n, 17)
    tb(sl, x+Cm(2.3), y+Cm(0.15), Cm(13), Cm(1.3), t, sz=22, b=True, c=CP)
    rect(sl, CA, x+Cm(2.3), y+Cm(1.4), Cm(3), Cm(0.06))
pgn(sl, 2)
 
# ========== Slide 3: Background ==========
sl = prs.slides.add_slide(blank)
tbar(sl, '一、研究背景与意义')
sdeco(sl)
tx = sl.shapes.add_textbox(Cm(1.5), Cm(3.3), Cm(31), Cm(14)); tf = tx.text_frame; tf.word_wrap = True
para(tf, 'IF钢与铜复合板', sz=20, b=True, c=CP)
blt(tf, 'IF钢（无间隙原子钢）：添加Ti、Nb固定C/N原子，具备优异的深冲性能和非时效性')
blt(tf, '铜：出色的导电性、导热性及良好的耐腐蚀性')
blt(tf, '复合后兼具高强度、易加工性和优良导电导热特性')
blt(tf, '应用领域：电子设备散热部件、汽车发动机关键组件等')
para(tf, '', sz=6)
para(tf, '研究必要性', sz=20, b=True, c=CP)
blt(tf, '快速加热工艺可显著影响复合板组织和性能，但保温时间优化研究严重不足')
blt(tf, '缺乏多温度梯度（830/850/870℃）下保温时间的系统研究')
blt(tf, '保温时间与界面结合、晶粒细化、力学性能的定量关系尚不明确')
blt(tf, '本研究填补该领域空白，为生产工艺优化提供可靠理论依据', b=True)
pgn(sl, 3)
 
# ========== Slide 4: Literature ==========
sl = prs.slides.add_slide(blank)
tbar(sl, '二、国内外研究现状')
sdeco(sl)
tx = sl.shapes.add_textbox(Cm(1.5), Cm(3.3), Cm(31), Cm(14)); tf = tx.text_frame; tf.word_wrap = True
para(tf, '异种金属复合界面扩散行为', sz=18, b=True, c=CS)
blt(tf, '钛/钢复合板：750~800℃扩散退火时界面剪切强度最高可达210MPa')
blt(tf, '过高/过低温度均会导致界面剪切强度下降，存在最佳温度窗口')
para(tf, '', sz=6)
para(tf, '热处理工艺对IF钢晶粒度的影响', sz=18, b=True, c=CS)
blt(tf, '加热速度、温度、保温时间和冷却速度均对晶粒度产生显著影响')
blt(tf, 'CSP线数据：碳含量50ppm以内时综合性能最优（r>2.0, n>0.22）')
para(tf, '', sz=6)
para(tf, '保温时间与界面剪切强度的相关性', sz=18, b=True, c=CS)
blt(tf, '适当延长保温时间促进元素扩散，增强界面结合')
blt(tf, '过长保温时间可能生成脆性相，导致强度下降')
blt(tf, '存在最佳保温时间窗口使界面性能最优')
pgn(sl, 4)
 
# ========== Slide 5: Content & Innovation ==========
sl = prs.slides.add_slide(blank)
tbar(sl, '三、研究内容与创新点')
sdeco(sl)
 
tx = sl.shapes.add_textbox(Cm(1.5), Cm(3.3), Cm(15), Cm(6)); tf = tx.text_frame; tf.word_wrap = True
para(tf, '研究内容', sz=20, b=True, c=CP)
blt(tf, '多温度梯度：830℃、850℃、870℃')
blt(tf, '保温时间：0s / 4s / 8s / 12s')
blt(tf, '金相+晶粒度+硬度 → 确定最优参数')
blt(tf, '拉伸试验 + 界面SEM断口检测')
 
tx2 = sl.shapes.add_textbox(Cm(17.5), Cm(3.3), Cm(15), Cm(6)); tf2 = tx2.text_frame; tf2.word_wrap = True
para(tf2, '创新点', sz=20, b=True, c=CP)
blt(tf2, '多温度梯度系统对比：构建温度-时间-组织性能完整图谱')
blt(tf2, '快速加热 vs 常规工艺对比：揭示快速加热对界面结合和晶粒细化的影响机制')
blt(tf2, '宏微观结合：宏观力学测试（拉伸/硬度）+ 微观表征（SEM断口分析）')
 
# Flow chart
para(tf, '', sz=4)
para(tf, '技术路线', sz=20, b=True, c=CP)
steps = [('准备试样', CH), ('到温入炉加热\n830/850/870℃', CH), ('保温\n0/4/8/12s', CH), ('水冷淬火', CH), ('性能测试\n金相/晶粒度/硬度', CG)]
for i, (st, sc) in enumerate(steps):
    x = Cm(1.5 + i*6.2)
    sh = rrect(sl, sc, x, Cm(11), Cm(5.5), Cm(1.5))
    sh.line.color.rgb = CS; sh.line.width = Pt(1)
    tb(sl, x+Cm(0.1), Cm(11)+Cm(0.1), Cm(5.3), Cm(1.3), st, sz=10, b=True, c=CP, a=PP_ALIGN.CENTER, anc=MSO_ANCHOR.MIDDLE)
    if i < len(steps)-1:
        tb(sl, x+Cm(5.7), Cm(11)+Cm(0.3), Cm(0.5), Cm(0.8), '>', sz=16, b=True, c=CS, a=PP_ALIGN.CENTER)
 
# TOC-like reminder at bottom
pgn(sl, 5)
 
# ========== Slide 6: Experimental ==========
sl = prs.slides.add_slide(blank)
tbar(sl, '四、实验材料与方法')
sdeco(sl)
 
tx = sl.shapes.add_textbox(Cm(1.5), Cm(3.3), Cm(15.5), Cm(5.5)); tf = tx.text_frame; tf.word_wrap = True
para(tf, '实验材料', sz=20, b=True, c=CP)
blt(tf, 'IF钢：0.002%C, Ti 0.04%, Nb 0.03%')
blt(tf, '纯铜：纯度 99.9%')
blt(tf, '复合板厚度：1.5mm')
para(tf, '', sz=4)
para(tf, '加热设备', sz=20, b=True, c=CP)
blt(tf, '常规电阻炉（硅碳棒加热）')
blt(tf, '温度范围：室温~1200℃，控温精度±1℃')
blt(tf, '升温速率：10℃/min')
 
chem_data = [
    ['元素', 'C', 'N', 'Ti', 'Nb', 'Si', 'Mn', 'Al', 'P', 'S'],
    ['含量(wt%)', chr(0x2264)+'0.003', chr(0x2264)+'0.003', '0.03~0.12', '0.01~0.06',
     chr(0x2264)+'0.03', chr(0x2264)+'0.15', '0.02~0.08', chr(0x2264)+'0.015', chr(0x2264)+'0.010'],
]
mtbl(sl, chem_data, Cm(1.5), Cm(9.5), Cm(15.5))
 
tx2 = sl.shapes.add_textbox(Cm(18), Cm(3.3), Cm(15), Cm(14)); tf2 = tx2.text_frame; tf2.word_wrap = True
para(tf2, '实验方案', sz=20, b=True, c=CP)
blt(tf2, '加热温度：830℃ / 850℃ / 870℃')
blt(tf2, '保温时间：0s / 4s / 8s / 12s')
blt(tf2, '冷却方式：水冷淬火（20℃去离子水）')
blt(tf2, '到温入炉，电子计时器精确计时')
para(tf2, '', sz=4)
para(tf2, '性能测试', sz=20, b=True, c=CP)
blt(tf2, '金相观察：蔡司Axioscope5（500×）')
blt(tf2, '晶粒度评级：ASTM E112 截点法')
blt(tf2, '显微硬度：HV-1000（25g/5s）')
blt(tf2, '拉伸试验：CMT5105（1mm/min）')
blt(tf2, 'SEM断口：喷金处理')
pgn(sl, 6)
 
# ========== Slide 7: Metallography & Grain Size ==========
sl = prs.slides.add_slide(blank)
tbar(sl, '五、结果与讨论——金相组织与晶粒度')
 
tx = sl.shapes.add_textbox(Cm(1.5), Cm(3.3), Cm(14), Cm(5)); tf = tx.text_frame; tf.word_wrap = True
para(tf, '金相组织观察', sz=20, b=True, c=CP)
blt(tf, '830℃：4s组织细小均匀(19.66um) → 8s最细(18.43um) → 12s略粗(21.78um)')
blt(tf, '850℃：随保温时间延长，晶粒持续长大(22.29→24.81→26.85um)')
blt(tf, '870℃：晶粒长大更显著(25.82→26.83→27.25um)，高温粗化明显')
 
grain_data = [
    ['条件', '晶粒尺寸范围(um)', '平均尺寸(um)', '晶粒度', '备注'],
    ['原件(未处理)', '8.69~38.94', '21.22', '8.6级', '原始状态'],
    ['830C-4s', '8.13~33.23', '19.66', '8.7级', chr(9733)+' 晶粒细化'],
    ['830C-8s', '9.64~31.29', '18.43', '8.8级', chr(9733)+chr(9733)+' 最细'],
    ['830C-12s', '9.19~31.37', '21.78', '8.6级', ''],
    ['850C-4s', '10~36.63', '22.29', '8.6级', ''],
    ['850C-8s', '13.1~37.19', '24.81', '8.4级', chr(8594)+' 综合优化'],
    ['850C-12s', '12.16~47.06', '26.85', '8.3级', ''],
    ['870C-4s', '10.21~41.98', '25.82', '8.3级', chr(8594)+' 塑性最优'],
    ['870C-8s', '14.19~42.14', '26.83', '8.3级', ''],
    ['870C-12s', '12.43~51.61', '27.25', '8.2级', '晶粒粗化'],
]
mtbl(sl, grain_data, Cm(1.5), Cm(8.5), Cm(31), cw=[Cm(5.5), Cm(8), Cm(5.5), Cm(4), Cm(5.5)])
pgn(sl, 7)
 
# ========== Slide 8: Hardness ==========
sl = prs.slides.add_slide(blank)
tbar(sl, '五、结果与讨论——硬度分析')
 
tx = sl.shapes.add_textbox(Cm(1.5), Cm(3.3), Cm(31), Cm(7)); tf = tx.text_frame; tf.word_wrap = True
para(tf, '不同温度下硬度随保温时间的变化趋势', sz=20, b=True, c=CP)
blt(tf, '总体趋势：随保温时间延长，硬度总体呈下降趋势', b=True)
blt(tf, '830℃： 4s最高 139.62HV -> 8s 134.08HV -> 12s 132.75HV')
blt(tf, '850℃： 4s 137.17HV -> 8s 138.69HV(略升) -> 12s 131.33HV(显著降)')
blt(tf, '870℃： 4s最高 138.45HV -> 8s 131.89HV -> 12s 130.80HV')
para(tf, '', sz=4)
para(tf, '晶粒度与硬度的响应规律（Hall-Petch关系）', sz=18, b=True, c=CP)
blt(tf, '晶粒细化强化：热处理后硬度普遍高于原始(130.52HV)，再结晶形成细小等轴晶')
blt(tf, '晶粒长大软化：同一温度下随保温时间延长，晶粒增大 -> 硬度下降')
blt(tf, '830C晶粒细化显著 -> 硬度升高；870C晶粒长大加剧 -> 硬度降低')
 
hdata = [
    ['温度', '保温时间', '硬度平均值(HV)', '变化趋势'],
    ['原件', '-', '130.52', '基准'],
    ['830C', '4s', '139.62', chr(8593)+' 最高'],
    ['830C', '8s', '134.08', chr(8595)],
    ['830C', '12s', '132.75', chr(8595)],
    ['850C', '4s', '137.16', chr(8593)],
    ['850C', '8s', '138.69', chr(8593)+' 略升'],
    ['850C', '12s', '131.32', chr(8595)+chr(8595)+' 显著'],
    ['870C', '4s', '138.45', chr(8593)+' 最高'],
    ['870C', '8s', '131.89', chr(8595)],
    ['870C', '12s', '130.80', chr(8595)],
]
mtbl(sl, hdata, Cm(1.5), Cm(11.5), Cm(31), cw=[Cm(7), Cm(7), Cm(9), Cm(8)])
pgn(sl, 8)
 
# ========== Slide 9: Tensile ==========
sl = prs.slides.add_slide(blank)
tbar(sl, '五、结果与讨论——拉伸性能')
 
tx = sl.shapes.add_textbox(Cm(1.5), Cm(3.3), Cm(31), Cm(14)); tf = tx.text_frame; tf.word_wrap = True
para(tf, '最优参数下的拉伸试验结果', sz=20, b=True, c=CP)
para(tf, '', sz=4)
 
tdata = [
    ['参数', '试样', '抗拉强度(MPa)', '弹性模量(MPa)', '下屈服强度(MPa)', '伸长率(%)'],
    ['830C/4s', '第1组', '300', '73004.61', '150', '49.5'],
    ['', '第2组', '310', '78179.65', '107', '47'],
    ['850C/8s', '第1组', '305', '79596.25', '152', '64'],
    ['', '第2组', '305', '80177.94', '139', '60.5'],
    ['870C/4s', '第1组', '305', '76933.97', '154', '65.5'],
    ['', '第2组', '305', '73980.99', '156', '62.5'],
]
mtbl(sl, tdata, Cm(1.5), Cm(5.5), Cm(31), cw=[Cm(4.5), Cm(3.5), Cm(6), Cm(6.5), Cm(6), Cm(4.5)])
 
para(tf, '', sz=4)
para(tf, '拉伸性能分析', sz=18, b=True, c=CP)
blt(tf, '830C保温4s：强度较高(300~310MPa)，但伸长率较低(~48%)，塑性相对不足')
blt(tf, '850C保温8s：强度适中(305MPa)，伸长率显著提高(~62%)，综合性能优异', b=True)
blt(tf, '870C保温4s：强度与850C相近，伸长率最高(~64%)，塑性最优', b=True)
pgn(sl, 9)
 
# ========== Slide 10: SEM ==========
sl = prs.slides.add_slide(blank)
tbar(sl, '五、结果与讨论——界面断口SEM分析')
 
tx = sl.shapes.add_textbox(Cm(1.5), Cm(3.3), Cm(31), Cm(14)); tf = tx.text_frame; tf.word_wrap = True
para(tf, '830C保温4s 断口特征', sz=18, b=True, c=CS)
blt(tf, '可见大量韧窝（韧性断裂特征），存在撕裂棱（含脆性成分）')
blt(tf, '存在明显界面分离现象 -> IF钢与铜结合强度较弱')
blt(tf, '断裂机制：微孔洞形核 -> 长大聚合 -> 韧性断裂+界面脱粘')
para(tf, '', sz=4)
para(tf, '850C保温8s 断口特征（最佳）', sz=18, b=True, c=CS)
blt(tf, '韧窝尺寸更大、形状更规则、分布更均匀 -> 塑性更好')
blt(tf, '撕裂棱明显减少，界面分离现象显著减少')
blt(tf, '断裂机制以韧性断裂为主 -> 界面结合强度提高')
para(tf, '', sz=4)
para(tf, '870C保温4s 断口特征', sz=18, b=True, c=CS)
blt(tf, '韧窝尺寸较大但分布不均，存在较大孔洞')
blt(tf, '界面分离现象略有增加，脆性成分有所增加')
blt(tf, '高温导致界面扩散加剧，但结合强度有所下降')
para(tf, '', sz=6)
blt(tf, '小结：850C/8s 界面结合质量最佳，断裂以韧性断裂为主', b=True)
pgn(sl, 10)
 
# ========== Slide 11: Optimization ==========
sl = prs.slides.add_slide(blank)
tbar(sl, '五、结果与讨论——最优工艺参数综合对比')
 
tx = sl.shapes.add_textbox(Cm(1.5), Cm(3.3), Cm(31), Cm(14)); tf = tx.text_frame; tf.word_wrap = True
para(tf, '三组最优参数综合对比', sz=20, b=True, c=CP)
para(tf, '', sz=4)
 
odata = [
    ['参数', '晶粒度', '硬度(HV)', '抗拉强度(MPa)', '伸长率(%)', '综合评价'],
    ['830C/4s', '8.7级', '139.62', '300~310', '48.25', '硬度高，塑性一般'],
    ['850C/8s', '8.4级', '138.69', '305', '62.25', chr(9733)+' 综合性能优异'],
    ['870C/4s', '8.3级', '138.45', '305', '64.00', chr(9733)+' 塑性最优'],
]
mtbl(sl, odata, Cm(1.5), Cm(5.5), Cm(31), cw=[Cm(4.5), Cm(3.5), Cm(4.5), Cm(5.5), Cm(4.5), Cm(7.5)])
 
para(tf, '', sz=6)
para(tf, '与常规退火工艺对比', sz=18, b=True, c=CP)
blt(tf, '常规退火工艺：加热速度慢，保温时间长（数分钟到数十分钟）')
blt(tf, '常规工艺下：抗拉强度约400MPa，延伸率约25%，平均晶粒尺寸约15um')
blt(tf, '本工艺（快速加热+优化保温）：强度约305MPa，延伸率高达64%，界面结合良好')
para(tf, '', sz=4)
blt(tf, '结论：快速加热工艺在优化保温时间后，能够显著改善IF复合板的综合力学性能', b=True)
pgn(sl, 11)
 
# ========== Slide 12: Conclusion ==========
sl = prs.slides.add_slide(blank)
tbar(sl, '六、研究结论')
 
tb(sl, Cm(1.5), Cm(3.5), Cm(31), Cm(1.5), '主要研究结论', sz=24, b=True, c=CP)
 
cons = [
    ('1', '建立了保温时间-温度-组织性能响应模型',
     ['830C：晶粒细化显著，4s时平均晶粒尺寸19.66um，硬度最高139.62HV',
      '850C：晶粒长大加速，8s时综合性能最优（伸长率62.25%，硬度138.69HV）',
      '870C：晶粒粗化明显，4s时塑性最优（伸长率64%）']),
    ('2', '确定了三组最优工艺参数',
     ['830C保温4s | 850C保温8s | 870C保温4s',
      '三者均可实现良好界面结合和优异力学性能']),
    ('3', '界面分析验证了良好的结合质量',
     ['拉伸后界面未出现脱粘或分层现象',
      'SEM显示以韧性断裂为主，证实界面结合质量良好']),
]
for i, (num, title, items) in enumerate(cons):
    yb = Cm(5.8 + i*3.6)
    oval(sl, Cm(1.5), yb, Cm(1.2), CP, num, 16)
    tb(sl, Cm(3.2), yb+Cm(0.05), Cm(29), Cm(1.1), title, sz=17, b=True, c=CP)
    for j, item in enumerate(items):
        tb(sl, Cm(3.5), yb+Cm(1.2+j*0.7), Cm(28), Cm(0.7), chr(8226)+' '+item, sz=14, c=CD)
pgn(sl, 12)
 
# ========== Slide 13: Outlook ==========
sl = prs.slides.add_slide(blank)
tbar(sl, '七、研究展望')
sdeco(sl)
 
tx = sl.shapes.add_textbox(Cm(1.5), Cm(3.5), Cm(31), Cm(14)); tf = tx.text_frame; tf.word_wrap = True
para(tf, '长期服役性能评估', sz=18, b=True, c=CP)
blt(tf, '开展疲劳寿命试验，研究循环载荷下的裂纹萌生和扩展机制')
blt(tf, '进行耐腐蚀性能研究，评估不同腐蚀介质中的服役行为')
para(tf, '', sz=6)
para(tf, '工业化应用推广', sz=18, b=True, c=CP)
blt(tf, '开发适合工业化生产的快速加热设备')
blt(tf, '优化生产工艺，提高生产效率和产品质量稳定性')
blt(tf, '开展中试试验，验证工艺在工业化生产中的可行性')
para(tf, '', sz=6)
para(tf, '跨学科交叉融合', sz=18, b=True, c=CP)
blt(tf, '材料科学 + 机械工程 + 自动化控制协同创新')
blt(tf, '推动快速加热工艺的规模化应用')
pgn(sl, 13)
 
# ========== Slide 14: Thanks ==========
sl = prs.slides.add_slide(blank)
rect(sl, CP, 0, 0, W, H)
rect(sl, CA, 0, 0, W, Cm(0.3))
tb(sl, Cm(2), Cm(3.5), Cm(30), Cm(2.5), '致  谢', sz=48, b=True, c=CW, a=PP_ALIGN.CENTER)
rect(sl, CG, Cm(13), Cm(6.5), Cm(8), Cm(0.06))
tb(sl, Cm(2), Cm(8), Cm(30), Cm(3),
   '感谢指导教师唐恩老师的悉心指导\n感谢实验室师兄师姐们的无私帮助\n感谢浙江水利水电学院机械工程学院的培养',
   sz=19, c=CH, a=PP_ALIGN.CENTER)
tb(sl, Cm(2), Cm(12.5), Cm(30), Cm(2), '请各位评委老师批评指正！',
   sz=28, b=True, c=CG, a=PP_ALIGN.CENTER)
tb(sl, Cm(2), H-Cm(1.5), Cm(30), Cm(1), '张新  |  机械工程学院  |  材料成型及控制工程  |  2025届',
   sz=12, c=CMG, a=PP_ALIGN.CENTER)
pgn(sl, 14)
 
# === SAVE ===
out = r'D:\桌面\abc\张新_毕业设计答辩PPT.pptx'
prs.save(out)
print(f'PPT saved: {out}')
print(f'Slides: {len(prs.slides)}')
